"""
Document parsing service.

Extracts real text and tables from uploaded files. No mocked output --
every branch below actually opens and reads the file with the relevant
library.
"""
from __future__ import annotations

import csv
import io
import os
from typing import Tuple, List

from pypdf import PdfReader
from pptx import Presentation
import openpyxl


SUPPORTED_EXTENSIONS = {".pdf", ".pptx", ".xlsx", ".xlsm", ".csv", ".txt", ".md"}


def detect_file_type(filename: str) -> str:
    ext = os.path.splitext(filename)[1].lower()
    return {
        ".pdf": "pdf",
        ".pptx": "pptx",
        ".xlsx": "xlsx",
        ".xlsm": "xlsx",
        ".csv": "csv",
        ".txt": "text",
        ".md": "markdown",
    }.get(ext, "unknown")


def parse_pdf(file_path: str) -> Tuple[str, List[List[List[str]]]]:
    """Returns (full_text, tables). PDF table extraction here is heuristic:
    pypdf doesn't do true table detection, so we extract text per page;
    real table structure recognition is layered on for xlsx/csv sources."""
    reader = PdfReader(file_path)
    pages_text = []
    for i, page in enumerate(reader.pages):
        text = page.extract_text() or ""
        pages_text.append(f"\n--- Page {i + 1} ---\n{text}")
    return "".join(pages_text), []


def parse_pptx(file_path: str) -> Tuple[str, List[List[List[str]]]]:
    prs = Presentation(file_path)
    slides_text = []
    tables: List[List[List[str]]] = []
    for i, slide in enumerate(prs.slides):
        slide_lines = [f"\n--- Slide {i + 1} ---"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs)
                    if line.strip():
                        slide_lines.append(line)
            if shape.has_table:
                table_data = []
                for row in shape.table.rows:
                    table_data.append([cell.text for cell in row.cells])
                tables.append(table_data)
            if shape.shape_type == 13:  # PICTURE
                slide_lines.append("[image on slide]")
        slides_text.append("\n".join(slide_lines))
    return "\n".join(slides_text), tables


def parse_xlsx(file_path: str) -> Tuple[str, List[List[List[str]]]]:
    wb = openpyxl.load_workbook(file_path, data_only=True)
    all_tables: List[List[List[str]]] = []
    text_parts = []
    for sheet in wb.worksheets:
        rows = []
        for row in sheet.iter_rows(values_only=True):
            rows.append(["" if c is None else str(c) for c in row])
        # Trim fully empty trailing rows
        while rows and all(cell == "" for cell in rows[-1]):
            rows.pop()
        if rows:
            all_tables.append(rows)
            text_parts.append(f"\n--- Sheet: {sheet.title} ---\n")
            for row in rows[:200]:  # cap for very large sheets
                text_parts.append(" | ".join(row))
    return "\n".join(text_parts), all_tables


def parse_csv(file_path: str) -> Tuple[str, List[List[List[str]]]]:
    with open(file_path, newline="", encoding="utf-8", errors="ignore") as f:
        reader = csv.reader(f)
        rows = [row for row in reader]
    text = "\n".join(" | ".join(row) for row in rows[:500])
    return text, [rows] if rows else []


def parse_text(file_path: str) -> Tuple[str, List[List[List[str]]]]:
    with open(file_path, encoding="utf-8", errors="ignore") as f:
        return f.read(), []


def parse_document(file_path: str, file_type: str) -> Tuple[str, List[List[List[str]]]]:
    """
    Dispatches to the correct real parser. Returns (extracted_text, tables).
    Raises ValueError for unsupported types -- callers should catch this and
    mark the document as failed rather than silently faking output.
    """
    if file_type == "pdf":
        return parse_pdf(file_path)
    if file_type == "pptx":
        return parse_pptx(file_path)
    if file_type == "xlsx":
        return parse_xlsx(file_path)
    if file_type == "csv":
        return parse_csv(file_path)
    if file_type in ("text", "markdown"):
        return parse_text(file_path)
    raise ValueError(f"Unsupported file type: {file_type}")


def chunk_text(text: str, chunk_size: int = 800, overlap: int = 120) -> List[str]:
    """
    Simple sliding-window chunker measured in characters (fast, dependency-free).
    Skips empty chunks. Overlap preserves context across chunk boundaries for
    retrieval quality.
    """
    text = text.strip()
    if not text:
        return []
    chunks = []
    start = 0
    n = len(text)
    while start < n:
        end = min(start + chunk_size, n)
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
        if end == n:
            break
        start = end - overlap
    return chunks
